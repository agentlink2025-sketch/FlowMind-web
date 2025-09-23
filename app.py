from flask import Flask, request, jsonify, Response, send_from_directory, session, redirect, url_for
from flask_cors import CORS
from zhipuai import ZhipuAI
import json
import requests
from requests.exceptions import Timeout
import time
import httpx
from openai import OpenAI
from langchain_core.prompts import ChatPromptTemplate
import os
import pandas as pd
import uuid
from werkzeug.utils import secure_filename
from apscheduler.schedulers.background import BackgroundScheduler
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt
import tempfile
from ai_ppt import AIPPT
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import re
import subprocess

app = Flask(__name__)

# 配置 CORS，允许跨域请求并支持凭证
CORS(app, supports_credentials=True)

# MySQL数据库配置（用于问答记录）
MYSQL_CONFIG = {
    'host': os.getenv('MYSQL_HOST', 'qa-db-mysql.ns-rlnc5x3h.svc'),
    'user': os.getenv('MYSQL_USER', 'root'),
    'password': os.getenv('MYSQL_PASSWORD', 'lmlqxrz9'),
    'database': os.getenv('MYSQL_DATABASE', 'qa_db')
}


# 构建MySQL数据库URI
MYSQL_DATABASE_URI = f"mysql+pymysql://{MYSQL_CONFIG['user']}:{MYSQL_CONFIG['password']}@{MYSQL_CONFIG['host']}/{MYSQL_CONFIG['database']}"

# 配置 Flask 应用
app.config.update(
    SECRET_KEY=os.urandom(24),
    SESSION_TYPE='filesystem',
    PERMANENT_SESSION_LIFETIME=timedelta(days=1),
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SECURE=True,
    SESSION_COOKIE_SAMESITE='Lax',
    # 用户认证使用SQLite
    SQLALCHEMY_DATABASE_URI='sqlite:///main.db',
    SQLALCHEMY_BINDS={
        'users': 'sqlite:///main.db',
        'chat_records': 'sqlite:///main.db',
        'qa_records': MYSQL_DATABASE_URI  # 问答记录使用MySQL
    },
    SQLALCHEMY_TRACK_MODIFICATIONS=False,
    UPLOAD_FOLDER='./uploads',
    MAX_CONTENT_LENGTH=100 * 1024 * 1024  # 100MB
)

# 初始化数据库
db = SQLAlchemy(app)

# 初始化登录管理器
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'index'
login_manager.session_protection = 'strong'

@login_manager.user_loader
def load_user(user_id):
    try:
        return User.query.get(int(user_id))
    except Exception as e:
        print(f"Error loading user: {str(e)}")
        return None

@login_manager.unauthorized_handler
def unauthorized():
    return jsonify({'error': '请先登录'}), 401

# 用户模型 - 使用SQLite
class User(UserMixin, db.Model):
    __bind_key__ = 'users'
    __tablename__ = 'user'
    id = db.Column(db.Integer, primary_key=True)
    phone = db.Column(db.String(20), unique=True, nullable=False)
    password_hash = db.Column(db.String(128))
    daily_queries = db.Column(db.Integer, default=0)
    last_query_date = db.Column(db.Date, default=datetime.now().date())
    
    def set_password(self, password):
        self.password_hash = generate_password_hash(password)
        
    def check_password(self, password):
        return check_password_hash(self.password_hash, password)
    
    def reset_daily_queries(self):
        try:
            if self.last_query_date != datetime.now().date():
                self.daily_queries = 0
                self.last_query_date = datetime.now().date()
                db.session.commit()
        except Exception as e:
            print(f"Error in reset_daily_queries: {str(e)}")
            db.session.rollback()
    
    def can_make_query(self):
        try:
            self.reset_daily_queries()
            return self.daily_queries < 100
        except Exception as e:
            print(f"Error in can_make_query: {str(e)}")
            return False
    
    def increment_queries(self):
        try:
            self.reset_daily_queries()
            self.daily_queries += 1
            db.session.add(self)
            db.session.commit()
            return True
        except Exception as e:
            print(f"Error in increment_queries: {str(e)}")
            db.session.rollback()
            return False
            
    def get_id(self):
        try:
            return str(self.id)
        except Exception as e:
            print(f"Error in get_id: {str(e)}")
            return None

# 聊天记录模型 - 使用SQLite
class ChatRecord(db.Model):
    __bind_key__ = 'users'
    __tablename__ = 'chat_record'

    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    question = db.Column(db.Text, nullable=False)
    question_time = db.Column(db.DateTime, default=datetime.now)
    intent_result = db.Column(db.Text)
    model_response = db.Column(db.Text)
    file_id = db.Column(db.String(36))

    user = db.relationship('User', backref=db.backref('chat_records', lazy=True))
    
    def to_dict(self):
        return {
            'id': self.id,
            'user_id': self.user_id,
            'question': self.question,
            'question_time': self.question_time.isoformat(),
            'intent_result': json.loads(self.intent_result) if self.intent_result else None,
            'model_response': self.model_response,
            'file_id': self.file_id
        }

# 问答记录模型 - 使用MySQL
class QARecord(db.Model):
    __bind_key__ = 'qa_records'  # 使用MySQL数据库
    __tablename__ = 'qa_record'

    id = db.Column(db.Integer, primary_key=True)
    conversation_id = db.Column(db.String(20), nullable=False, index=True)
    user_id = db.Column(db.Integer, nullable=False)  # 不设置外键，因为跨数据库
    question = db.Column(db.Text, nullable=False)
    answer = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.now)
    user_tag = db.Column(db.String(50))
    answer_quality = db.Column(db.String(20))
    meta_data = db.Column(db.JSON)

    def to_dict(self):
        return {
            'id': self.id,
            'conversation_id': self.conversation_id,
            'user_id': self.user_id,
            'question': self.question,
            'answer': self.answer,
            'timestamp': self.timestamp.isoformat(),
            'user_tag': self.user_tag,
            'answer_quality': self.answer_quality,
            'meta_data': self.meta_data
        }

    @staticmethod
    def generate_conversation_id():
        """生成对话ID，格式：YYYYMMDD_XXX"""
        today = datetime.now().strftime('%Y%m%d')
        max_record = QARecord.query.filter(
            QARecord.conversation_id.like(f'{today}_%')
        ).order_by(QARecord.conversation_id.desc()).first()
        
        if max_record:
            last_num = int(max_record.conversation_id.split('_')[1])
            new_num = last_num + 1
        else:
            new_num = 1
            
        return f"{today}_{new_num:03d}"

# MySQL中的用户模型
class MySQLUser(db.Model):
    __bind_key__ = 'qa_records'
    __tablename__ = 'user'
    id = db.Column(db.Integer, primary_key=True)
    phone = db.Column(db.String(20), unique=True, nullable=False)
    password_hash = db.Column(db.String(128))
    daily_queries = db.Column(db.Integer, default=0)
    last_query_date = db.Column(db.Date, default=datetime.now().date())

    def reset_daily_queries(self):
        try:
            if self.last_query_date != datetime.now().date():
                self.daily_queries = 0
                self.last_query_date = datetime.now().date()
                db.session.commit()
        except Exception as e:
            print(f"Error in MySQL reset_daily_queries: {str(e)}")
            db.session.rollback()
    
    def can_make_query(self):
        try:
            self.reset_daily_queries()
            return self.daily_queries < 100
        except Exception as e:
            print(f"Error in MySQL can_make_query: {str(e)}")
            return False
    
    def increment_queries(self):
        try:
            self.reset_daily_queries()
            self.daily_queries += 1
            db.session.add(self)
            db.session.commit()
            return True
        except Exception as e:
            print(f"Error in MySQL increment_queries: {str(e)}")
            db.session.rollback()
            return False

# 用户每日使用记录模型 - 使用MySQL
class UserDailyUsage(db.Model):
    __bind_key__ = 'qa_records'
    __tablename__ = 'user_daily_usage'
    
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False, index=True)
    date = db.Column(db.Date, nullable=False, index=True)
    usage_count = db.Column(db.Integer, default=0)
    created_at = db.Column(db.DateTime, default=datetime.now)
    updated_at = db.Column(db.DateTime, default=datetime.now, onupdate=datetime.now)
    
    __table_args__ = (
        db.UniqueConstraint('user_id', 'date', name='uix_user_date'),
    )
    
    @classmethod
    def get_or_create(cls, user_id, date=None):
        if date is None:
            date = datetime.now().date()
            
        usage = cls.query.filter_by(user_id=user_id, date=date).first()
        if not usage:
            usage = cls(user_id=user_id, date=date, usage_count=0)
            db.session.add(usage)
            db.session.commit()
        return usage
    
    def increment_usage(self):
        try:
            self.usage_count += 1
            self.updated_at = datetime.now()
            db.session.add(self)
            db.session.commit()
            return True
        except Exception as e:
            print(f"Error incrementing usage: {str(e)}")
            db.session.rollback()
            return False
    
    def can_make_query(self):
        return self.usage_count < 100

@app.route('/register', methods=['POST'])
def register():
    data = request.json
    phone = data.get('phone')
    password = data.get('password')
    
    if not phone or not password:
        return jsonify({'error': '手机号和密码不能为空'}), 400
    
    # 验证手机号格式
    if not phone.isdigit() or len(phone) != 11:
        return jsonify({'error': '无效的手机号格式'}), 400
    
    # 检查用户是否已存在（同时检查SQLite和MySQL）
    if User.query.filter_by(phone=phone).first() or MySQLUser.query.filter_by(phone=phone).first():
        return jsonify({'error': '该手机号已注册'}), 400
    
    # 创建新用户
    user = User(phone=phone)
    user.set_password(password)
    mysql_user = MySQLUser(phone=phone, password_hash=user.password_hash, daily_queries=0, last_query_date=datetime.now().date())
    try:
        db.session.add(user)
        db.session.flush()  # 先写入SQLite，获取id
        mysql_user.id = user.id  # 保持两边id一致
        db.session.add(mysql_user)
        db.session.commit()
        return jsonify({'message': '注册成功'}), 201
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': '注册失败'}), 500

@app.route('/login', methods=['POST'])
def login():
    try:
        data = request.json
        phone = data.get('phone')
        password = data.get('password')
        
        if not phone or not password:
            return jsonify({'error': '手机号和密码不能为空'}), 400
        
        user = User.query.filter_by(phone=phone).first()
        if user and user.check_password(password):
            # 登录用户
            login_user(user, remember=True)
            # 设置会话为永久
            session.permanent = True
            # 返回用户信息
            return jsonify({
                'message': '登录成功',
                'user': {
                    'phone': user.phone,
                    'daily_queries': user.daily_queries,
                    'remaining_queries': 100 - user.daily_queries
                }
            }), 200
        
        return jsonify({'error': '手机号或密码错误'}), 401
    except Exception as e:
        print(f"Error in login: {str(e)}")
        return jsonify({'error': '登录失败，请重试'}), 500

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return jsonify({'message': '已退出登录'}), 200

@app.route('/user/status')
@login_required
def user_status():
    user = current_user
    today = datetime.now().date()
    daily_usage = UserDailyUsage.get_or_create(user.id, today)
    
    return jsonify({
        'phone': user.phone,
        'daily_queries': daily_usage.usage_count,
        'remaining_queries': 100 - daily_usage.usage_count
    }), 200

# API Keys
DEEPSEEK_API_KEY = "sk-1c679cccfc00441fabc54b78a69e8b9c"
API_URL = "https://api.deepseek.com/v1/chat/completions"
INTENT_MODEL_URL = 'http://172.17.0.3:10052/v1'
INTENT_MODEL_URL = 'http://qjq-n7527ixuk1ttjbfq5-fzbecewvq-custom.service.onethingrobot.com/v1'

# 添加AiPPT API相关配置
AIPPT_APP_ID = "XXXXXXXX"  # 请替换为实际的APP ID
AIPPT_API_SECRET = "XXXXXXXXXXXXXXXXXXXXXXXX"  # 请替换为实际的API Secret
AIPPT_TEMPLATE_ID = "20240718489569D"  # PPT模板ID

# 初始化智谱AI客户端
#client = ZhipuAI(api_key="4701d076f5fc5b79919d7ed7d1fd220e.cZGekQXg0fypWIkY") 
#a78b2a8931a445e98f12a128fdc651ec.aMDS3pYc1xjpVcMw # 请填写您自己的APIKey
client = ZhipuAI(api_key="c6e5f8c041734d3a898697f83fd0bf44.PmZqIob2UELuVhGq")

file_storage = {}
scheduler = BackgroundScheduler()

# 文件清理任务
def clean_files():
    now = datetime.now()
    expired = []
    for fid, data in file_storage.items():
        if (now - datetime.fromisoformat(data['upload_time'])).seconds > 3600:
            expired.append(fid)
    for fid in expired:
        try:
            os.remove(file_storage[fid]['path'])
            del file_storage[fid]
        except: pass

scheduler.add_job(clean_files, 'interval', minutes=60)
scheduler.start()

# 辅助函数
def analyze_data(df, question):
    """生成数据分析提示"""
    columns = df.columns.tolist()
    sample = df.head(3).to_dict()
    return f"""你是一个专业数据分析师，请严格根据以下数据回答问题：
    
数据集特征：{columns}
数据样例：{sample}

用户问题：{question}

回答要求：
1. 基于数据实际情况
2. 包含关键指标计算
3. 使用Markdown格式
4. 重要数据用**加粗**强调"""
# 企业微信机器人 Webhook
webhook_url = "https://qyapi.weixin.qq.com/cgi-bin/webhook/send?key=cae9ff6f-7125-4fe8-bec6-1d2aa0743231"

UPLOAD_FOLDER = './uploads'
ALLOWED_EXTENSIONS = {'xlsx', 'csv'}

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_model(openai_api_base):
    """
    模型url
    """
    print('openai_api_base:')
    print(openai_api_base)
    httpx_client = httpx.Client(http2=True, verify=False)
    openai_api_key = "EMPTY"
    openai_api_base = openai_api_base

    model = OpenAI(
        api_key=openai_api_key,
        base_url=openai_api_base,
        http_client=httpx_client
    )
    return model

def generate_story_with_custom_model(prompt, timeout=500):
    try:
        qw_model = get_model(INTENT_MODEL_URL)
        print("--qwen model--")
        print(qw_model)

        prompts = '''这是一个专业的故事生成助手，请根据用户需求生成高质量的故事。

<文本>
"{text}"

'''
        prompt_template = ChatPromptTemplate.from_template(prompts)
        new_prompt = prompt_template.invoke({'text': prompt}).to_string()[:7000]
        
        start_time = time.time()
        print("new_prompt test:")
        print(new_prompt)
        chat_response = qw_model.chat.completions.create(
            model='/root/LLaMA-Factory/qwen_model',
            messages=[
                {"role": "system", "content": 'you are a professional storyteller'},
                {"role": "user", "content": new_prompt}
            ],
            top_p=0.00000001,
            max_tokens=2000,
            temperature=0.1
        )
        
        # if time.time() - start_time > timeout:
        #     return None
        pass
        print("chat_response:")
        print(chat_response)
        print("--content--")
        print(chat_response.choices[0].message.content)
        return chat_response.choices[0].message.content
    except Exception as e:
        print(f"Custom model error: {e}")
        return None

def call_deepseek(prompt, timeout=300):
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "deepseek-chat",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.1
    }
    try:
        start_time = time.time()
        response = requests.post(API_URL, json=data, headers=headers, timeout=timeout)
        if time.time() - start_time > timeout:
            return None
        print("--deepseek response--")
        print(response)
        res = response.json()['choices'][0]['message']['content']
        return res
    except (Timeout, requests.exceptions.RequestException):
        return None

def evaluate_story(story):
    evaluation_prompt = f"""Please evaluate the following Patagonia story based on these criteria:
1. Authenticity and alignment with Patagonia's values
2. Writing style and quality
3. Story structure and engagement
4. CSR principles integration
5. Overall impact and message

Story to evaluate:
{story}

Please provide a detailed evaluation with scores for each criterion (1-10) and specific feedback for improvement."""

    return call_deepseek(evaluation_prompt)

def send_text_message(content, mentioned_mobile_list=None):
    data = {
        "msgtype": "text",
        "text": {
            "content": content,
            "mentioned_mobile_list": mentioned_mobile_list or []
        }
    }
    response = requests.post(webhook_url, json=data)
    return response.json()

def send_markdown_message(content):
    data = {
        "msgtype": "markdown",
        "markdown": {
            "content": content
        }
    }
    response = requests.post(webhook_url, json=data)
    return response.json()

def excel_to_text(file_path):
    # 读取Excel或CSV并转为文本
    if file_path.endswith('.csv'):
        try:
            # 首先尝试 UTF-8 编码
            df = pd.read_csv(file_path, encoding='utf-8')
        except UnicodeDecodeError:
            try:
                # 如果 UTF-8 失败，尝试 GBK 编码
                df = pd.read_csv(file_path, encoding='gbk')
            except UnicodeDecodeError:
                # 如果 GBK 也失败，尝试 GB18030 编码
                df = pd.read_csv(file_path, encoding='gb18030')
    else:
        try:
            # 尝试使用 openpyxl 引擎
            df = pd.read_excel(file_path, engine='openpyxl')
        except Exception as e:
            try:
                # 如果失败，尝试使用 xlrd 引擎
                df = pd.read_excel(file_path, engine='xlrd')
            except Exception as e:
                raise Exception(f"无法读取Excel文件: {str(e)}")
    text = df.to_string(index=False)
    return text

@app.route('/')
def index():
    return send_from_directory('static', 'index.html')

@app.route('/chat')
@login_required
def chat_page():
    try:
        if not current_user or not current_user.is_authenticated:
            return redirect(url_for('index'))
        return send_from_directory('static', 'chat.html')
    except Exception as e:
        print(f"Error in chat_page: {str(e)}")
        return redirect(url_for('index'))

@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    user_message = data.get('message', '').strip().lower()
    
    def generate():
        try:
            # 根据意图处理请求
            result = None
            if "分析意图" in intent_result.get("intents", []):
                result = handle_analysis_intent(user_message, file_data)
            elif "ROI预测意图" in intent_result.get("intents", []):
                result = handle_roi_intent(user_message, file_data)
            elif "执行意图" in intent_result.get("intents", []):
                result = handle_execution_intent(user_message)
            elif "内容生成意图" in intent_result.get("intents", []):
                result = handle_content_generation_intent(user_message)
            else:
                result = "抱歉，我无法理解您的意图。请尝试重新描述您的问题。"

            if not result:
                result = "抱歉，处理您的请求时出现错误。请稍后重试。"

            # 在应用上下文中执行数据库操作
            with app.app_context():
                # 增加用户查询计数
                if not user.increment_queries():
                    yield f"data: {json.dumps({'content': '更新查询次数失败，请重试'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return

                # 保存聊天记录
                try:
                    import json
                    model_response_str = result
                    if not isinstance(model_response_str, str):
                        try:
                            model_response_str = json.dumps(model_response_str, ensure_ascii=False)
                        except Exception:
                            model_response_str = str(model_response_str)
                    chat_record = ChatRecord(
                        user_id=user.id,
                        question=user_message,
                        intent_result=json.dumps(intent_result),
                        model_response=model_response_str,
                        file_id=file_id
                    )
                    db.session.add(chat_record)
                    db.session.commit()
                except Exception as e:
                    print(f"Error saving chat record: {str(e)}")
                    db.session.rollback()

            # 生成报告标题
            report_title = generate_report_title(user_message, result)
            
            # 生成PPT报告
            try:
                pptx_path = text_to_pptxs(result, report_title)
                if pptx_path:
                    upload_file_to_wecom(pptx_path, webhook_url)
            except Exception as e:
                print(f"生成PPT报告失败: {str(e)}")
            
            # 流式返回结果
            yield f"data: {json.dumps({'content': result})}\n\n"
            yield "data: [DONE]\n\n"
            
        except Exception as e:
            print(f"Error in generate: {str(e)}")
            yield f"data: {json.dumps({'content': '抱歉，处理您的请求时出现错误。请稍后重试。'})}\n\n"
            yield "data: [DONE]\n\n"

    return Response(generate(), mimetype='text/event-stream')
            
    # except Exception as e:
    #     print(f"Error in chat_handlerss: {str(e)}")
    #     return jsonify({'error': '处理请求时出现错误'}), 500

@app.route('/upload_excel', methods=['POST'])
def upload_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        # Excel转文本
        excel_text = excel_to_text(file_path)

        # 用大模型分析
        # 你可以直接用现有的分析函数，比如调用GLM-4或自研模型
        # 这里以GLM-4为例
        system_prompt = "请对以下Excel内容进行分析，总结主要信息并给出建议："
        response = client.chat.completions.create(
            model="glm-4",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": excel_text},
            ],
            stream=False,
        )
        result = response.choices[0].message.content

        # 推送到企业微信
        send_text_message(result)

        return jsonify({'result': result})
    else:
        return jsonify({'error': 'Invalid file type'}), 400

# 路由
@app.route('/upload_excelss', methods=['POST'])
def upload_excelss():
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未选择文件'}), 400
        
        file = request.files['file']
        if not file or file.filename == '':
            return jsonify({'error': '无效文件'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': '不支持的文件类型，请上传Excel文件(.xlsx)或CSV文件(.csv)'}), 400
        
        # 检查文件大小
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)  # 重置文件指针
        
        max_size = 100 * 1024 * 1024  # 100MB
        if file_size > max_size:
            return jsonify({'error': f'文件大小超过限制，请上传小于100MB的文件'}), 413
        
        # 保存文件
        file_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        try:
            file.save(filepath)
        except Exception as e:
            return jsonify({'error': f'文件保存失败: {str(e)}'}), 500
        
        # 读取元数据
        try:
            if filename.endswith('.csv'):
                # 尝试不同的编码方式读取CSV文件
                encodings = ['utf-8', 'gbk', 'gb18030']
                df = None
                for encoding in encodings:
                    try:
                        df = pd.read_csv(filepath, encoding=encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        print(f"Error reading CSV with {encoding}: {str(e)}")
                        continue
                
                if df is None:
                    raise Exception("无法读取CSV文件，请确保文件编码正确")
            else:
                try:
                    # 尝试使用 openpyxl 引擎
                    df = pd.read_excel(filepath, engine='openpyxl')
                except Exception as e:
                    try:
                        # 如果失败，尝试使用 xlrd 引擎
                        df = pd.read_excel(filepath, engine='xlrd')
                    except Exception as e:
                        raise Exception(f"无法读取Excel文件: {str(e)}")
        except Exception as e:
            os.remove(filepath)  # 删除无效文件
            return jsonify({'error': f'文件读取失败: {str(e)}'}), 400
        
        # 验证数据量
        if len(df) > 100000 or len(df.columns) > 50:
            os.remove(filepath)
            return jsonify({'error': '数据量过大，请确保数据行数不超过100,000行，列数不超过50列'}), 400
        
        # 存储元数据
        try:
            file_storage[file_id] = {
                'path': filepath,
                'columns': df.columns.tolist(),
                'sample': df.head(3).to_dict(),
                'upload_time': datetime.now().isoformat(),
                'size': f"{os.path.getsize(filepath)/1024:.1f}KB"
            }
        except Exception as e:
            os.remove(filepath)
            return jsonify({'error': f'元数据存储失败: {str(e)}'}), 500
        
        # 返回成功响应
        response_data = {
            'fileId': file_id,
            'fileName': filename,
            'fileSize': file_storage[file_id]['size'],
            'columns': file_storage[file_id]['columns']
        }
        
        return jsonify(response_data)
        
    except Exception as e:
        # 确保清理任何已创建的文件
        if 'filepath' in locals() and os.path.exists(filepath):
            os.remove(filepath)
        return jsonify({'error': f'文件处理失败: {str(e)}'}), 500

def recognize_intent(text):
    """识别用户意图，扩展GEO营销场景和广告参数生成意图，并支持IP位置查询意图"""
    prompt = f"""请仔细分析以下用户输入，识别其意图：\n\n用户输入：{text}\n\n可能的意图包括：\n1. 分析意图：用户想要分析数据，可能包含文件上传。典型场景包括：\n   - 指标计算\n   - 对比分析\n   - 数据统计\n   - 趋势分析\n   - 原因分析\n\n2. ROI预测意图：用户想要进行ROI预测分析，可能包含文件上传。典型场景包括：\n   - 给出优化方案\n   - ROI预测\n   - 效果评估\n   - 收益分析\n   - 投入产出比分析\n\n3. 执行意图：用户想要执行特定操作。典型场景包括：\n   - 根据最优方案进行执行\n   - 生成广告参数\n   - 广告投放参数\n   - 自动投放\n   - 配置调整\n   - 方案实施\n\n4. 广告参数生成意图：用户需要生成广告投放参数json。典型场景包括：\n   - 生成广告参数\n   - 广告投放参数\n\n5. 内容生成意图：用户需要生成特定内容。典型场景包括：\n   - 生成营销文案\n   - 生成产品描述\n   - 生成社交媒体内容\n   - 生成广告创意\n   - 生成品牌故事\n\n6. GEO营销-收录统计：用户想要查询某个关键词或品牌在Google等搜索引擎的收录情况。\n   - 典型场景：收录量、排名、曝光、长尾词收录等\n\n7. GEO营销-收录解决方案生成：用户想要获得提升收录的具体解决方案。\n   - 典型场景：SEO优化建议、内容调整、外链建设、技术方案等\n\n8. IP位置查询意图：用户想要查询某个IP地址的地理位置。\n   - 典型场景：IP归属地、IP位置、IP地理信息、IP查询\n\n用户可以同时有多个意图，比如同时需要分析和ROI预测。\n\n请严格按照以下JSON格式返回识别结果，不要包含任何其他内容：\n{{\n    \"intents\": [\"分析意图\", \"ROI预测意图\", \"执行意图\", \"广告参数生成意图\", \"内容生成意图\", \"GEO营销-收录统计\", \"GEO营销-收录解决方案生成\", \"IP位置查询意图\"],\n    \"has_file\": true,\n    \"details\": \"具体执行内容\"\n}}\n\n注意：\n1. 必须返回完整的JSON对象\n2. intents必须是数组，可以包含多个意图\n3. has_file必须是布尔值\n4. details必须是字符串\n5. 不要包含任何注释或额外说明\n"""
    try:
        response = client.chat.completions.create(
            model="glm-4",
            messages=[{"role": "user", "content": prompt}],
            stream=False
        )
        result = response.choices[0].message.content
        # 清理可能的额外内容
        result = result.strip()
        if result.startswith('```json'):
            result = result[7:]
        if result.endswith('```'):
            result = result[:-3]
        result = result.strip()
        # 处理BOM头
        if result and ord(result[0]) == 0xfeff:
            result = result[1:]
        # 解析JSON结果
        intent_data = json.loads(result)
        # 验证返回的数据结构
        if not isinstance(intent_data, dict):
            raise ValueError("返回结果不是JSON对象")
        if "intents" not in intent_data or not isinstance(intent_data["intents"], list):
            raise ValueError("intents字段格式错误")
        if "has_file" not in intent_data or not isinstance(intent_data["has_file"], bool):
            raise ValueError("has_file字段格式错误")
        if "details" not in intent_data or not isinstance(intent_data["details"], str):
            raise ValueError("details字段格式错误")
        print("--intent recognition result--")
        print(intent_data)
        return intent_data
    except json.JSONDecodeError as e:
        print(f"JSON解析错误: {e}")
        print(f"原始返回内容: {result}")
        return {"intents": ["分析意图"], "has_file": False, "details": ""}
    except Exception as e:
        print(f"意图识别错误: {e}")
        return {"intents": ["分析意图"], "has_file": False, "details": ""}

def handle_analysis_intent(text, file_data=None):
    """处理分析意图，实现结构化数据分析工作流"""
    try:
        if file_data:
            df = pd.read_excel(file_data['path'])
            
            # 检查是否是关键词分析请求
            is_keyword_analysis = any(keyword in text.lower() for keyword in [
                "关键词", "关键字", "搜索词", "搜索关键字", "搜索关键词",
                "推荐关键词", "关键词推荐", "关键词筛选"
            ])
            
            if is_keyword_analysis:
                # 关键词分析提示词
                keyword_prompt = f"""请对以下数据进行关键词分析：

数据特征：{df.columns.tolist()}
数据样例：{df.head(3).to_dict()}
数据形状：{df.shape}

用户需求：{text}

请提供以下分析：
1. 关键词筛选标准：
   - 搜索量
   - 竞争度
   - 转化率
   - 相关性
   - 商业价值

2. 推荐15个关键词，每个关键词包含：
   - 关键词名称
   - 搜索量（月均）
   - 竞争度（1-10）
   - 预计转化率
   - 商业价值评分（1-10）
   - 推荐理由

3. 关键词分类：
   - 高价值关键词（搜索量高、竞争度适中）
   - 长尾关键词（搜索量适中、竞争度低）
   - 蓝海关键词（搜索量增长快、竞争度低）

4. 投放建议：
   - 预算分配建议
   - 投放策略建议
   - 优化方向建议

要求：
1. 使用Markdown格式
2. 重要数据用**加粗**强调
3. 确保数据准确性和可操作性
4. 提供具体的优化建议
5. 包含数据支持
6. 结构清晰，层次分明"""
                
                keyword_analysis = call_deepseek(keyword_prompt, timeout=300)
                if not keyword_analysis:
                    raise Exception("关键词分析失败")
                return keyword_analysis
            
            # 1. 数据探索分析
            exploration_prompt = f"""请对以下数据进行探索性分析：

数据特征：{df.columns.tolist()}
数据样例：{df.head(3).to_dict()}
数据形状：{df.shape}
数据类型：{df.dtypes.to_dict()}
数值列统计：{df.describe().to_dict() if not df.empty else '无数据'}

请提供：
1. 数据质量评估
2. 关键指标识别
3. 数据分布特征
4. 异常值检测
5. 相关性分析建议

要求：
1. 所有数值必须完整展示，不能省略
2. 百分比必须精确到小数点后2位
3. 重要指标必须单独列出并加粗
4. 使用Markdown格式输出
5. 重要发现用**加粗**标记
6. 确保所有数据完整展示，不要使用省略号"""
            
            exploration_result = call_deepseek(exploration_prompt, timeout=300)
            if not exploration_result:
                raise Exception("数据探索分析失败")
            
            # 2. 问题分析
            analysis_prompt = f"""基于以下数据探索结果和用户问题，制定详细的分析方案：

数据探索结果：
{exploration_result}

用户问题：{text}

请提供：
1. 分析目标分解
2. 分析步骤规划
3. 所需分析方法
4. 预期输出结果
5. 潜在风险提示

要求：
1. 所有数值必须完整展示，不能省略
2. 百分比必须精确到小数点后2位
3. 重要指标必须单独列出并加粗
4. 使用Markdown格式输出
5. 重要内容用**加粗**标记
6. 确保所有数据完整展示，不要使用省略号"""
            
            analysis_plan = call_deepseek(analysis_prompt, timeout=300)
            if not analysis_plan:
                raise Exception("分析方案制定失败")
            
            # 3. 执行分析
            execution_prompt = f"""请根据以下分析方案执行具体分析：

分析方案：
{analysis_plan}

原始数据特征：
- 列名：{df.columns.tolist()}
- 数据样例：{df.head(3).to_dict()}
- 数据统计：{df.describe().to_dict() if not df.empty else '无数据'}

要求：
1. 严格按照分析方案执行
2. 提供详细的计算过程
3. 使用Markdown格式
4. 重要发现用**加粗**标记
5. 包含数据支持
6. 结构清晰，层次分明
7. 所有数值必须完整展示，不能省略
8. 百分比必须精确到小数点后2位
9. 重要指标必须单独列出并加粗
10. 确保所有数据完整展示，不要使用省略号
11. 对于关键指标（如高效区占比等），必须：
    - 单独列出并详细说明
    - 提供完整的数值
    - 解释其含义和影响
    - 与其他指标的关系
    - 建议的改进方向"""
            
            analysis_result = call_deepseek(execution_prompt, timeout=300)
            if not analysis_result:
                raise Exception("分析执行失败")
            
            # 4. 生成最终报告
            report_prompt = f"""请整合以下分析结果，生成完整的分析报告：

数据探索结果：
{exploration_result}

分析方案：
{analysis_plan}

分析执行结果：
{analysis_result}

要求：
1. 生成结构化的分析报告
2. 包含以下部分：
   - 执行摘要
   - 数据概览
   - 主要发现
   - 详细分析
   - 结论建议
3. 使用Markdown格式
4. 重要内容用**加粗**标记
5. 确保逻辑连贯
6. 突出关键发现
7. 所有数值必须完整展示，不能省略
8. 百分比必须精确到小数点后2位
9. 重要指标必须单独列出并加粗
10. 确保所有数据完整展示，不要使用省略号
11. 对于关键指标（如高效区占比等），必须：
    - 单独列出并详细说明
    - 提供完整的数值
    - 解释其含义和影响
    - 与其他指标的关系
    - 建议的改进方向"""
            
            final_report = call_deepseek(report_prompt, timeout=300)
            if not final_report:
                raise Exception("报告生成失败")
            
            return final_report
            
        else:
            # 无数据文件时的分析
            prompt = f"""请基于以下问题进行分析：

问题：{text}

要求：
1. 提供分析框架
2. 说明所需数据
3. 建议分析方法
4. 预期分析结果
5. 使用Markdown格式
6. 重要内容用**加粗**标记
7. 所有数值必须完整展示，不能省略
8. 百分比必须精确到小数点后2位
9. 重要指标必须单独列出并加粗
10. 确保所有数据完整展示，不要使用省略号"""
            
            result = call_deepseek(prompt, timeout=300)
            if not result:
                raise Exception("分析生成失败")
            return result
            
    except Exception as e:
        error_msg = f"数据分析失败: {str(e)}"
        send_markdown_message(error_msg)
        return error_msg

def handle_roi_intent(text, file_data=None):
    """处理ROI预测意图，统一用 deepseek API"""
    if file_data:
        df = pd.read_excel(file_data['path'])
        roi_prediction = predict_roi_with_model(df, text)
        prompt = f"""你是一个专业的ROI分析师，请基于以下数据生成多个优化方案并预测ROI：\n\n数据特征：{df.columns.tolist()}\n数据样例：{df.head(3).to_dict()}\n\nROI预测结果：\n{roi_prediction['report'] if roi_prediction else 'ROI预测失败'}\n\n问题：{text}\n\n要求：\n1. 生成3-5个不同的优化方案，每个方案必须包含：\n   - 方案背景和目标\n   - 具体实施步骤\n   - 所需资源清单\n   - 时间规划表\n   - 关键里程碑\n2. 对每个方案进行详细的ROI预测，包括：\n   ### 投入成本预测\n   - 直接成本：\n     * 人力成本（按岗位和工时计算）\n     * 设备/工具成本\n     * 软件/系统成本\n     * 培训成本\n     * 运营成本\n   - 间接成本：\n     * 管理成本\n     * 机会成本\n     * 风险准备金\n   - 成本分摊计划（按时间维度）\n   ### 收益预测\n   - 直接收益：\n     * 收入增长\n     * 成本节约\n     * 效率提升\n     * 质量改进\n   - 间接收益：\n     * 品牌价值提升\n     * 客户满意度提升\n     * 市场份额增长\n     * 竞争优势增强\n   - 收益时间分布\n   ### 财务指标预测\n   - 投资回收期（PP）\n   - 净现值（NPV）\n   - 内部收益率（IRR）\n   - 投资回报率（ROI）\n   - 盈亏平衡点分析\n   - 敏感性分析\n   ### 风险评估\n   - 市场风险\n   - 技术风险\n   - 执行风险\n   - 财务风险\n   - 风险应对策略\n3. 使用Markdown格式，结构如下：\n   ## 方案1：方案名称\n   ### 方案概述\n   - 背景和目标\n   - 实施范围\n   - 预期效果\n   ### 实施计划\n   - 具体步骤\n   - 时间安排\n   - 资源需求\n   - 里程碑节点\n   ### ROI预测分析\n   [上述所有ROI预测内容]\n   ### 方案评估\n   - 优势分析\n   - 劣势分析\n   - 机会分析\n   - 威胁分析\n   - 可行性评估\n   [其他方案类似]\n4. 最后给出最优方案推荐，包括：\n   ### 方案对比分析\n   - 各方案关键指标对比表\n   - 优劣势对比分析\n   - 风险对比分析\n   ### 最优方案选择\n   - 选择理由\n   - 实施建议\n   - 预期效果\n   - 风险控制措施\n   - 成功关键因素\n5. 重要数据用**加粗**强调，所有预测过程需要清晰展示\n"""
    else:
        prompt = f"""你是一个专业的ROI分析师，请基于以下问题构建多个优化方案并预测ROI：\n\n问题：{text}\n\n要求：\n1. 构建3-5个合理的假设场景作为优化方案，每个方案必须包含：\n   - 方案背景和目标\n   - 具体实施步骤\n   - 所需资源清单\n   - 时间规划表\n   - 关键里程碑\n2. 对每个方案进行详细的ROI预测，包括：\n   ### 投入成本预测\n   - 直接成本：\n     * 人力成本（按岗位和工时计算）\n     * 设备/工具成本\n     * 软件/系统成本\n     * 培训成本\n     * 运营成本\n   - 间接成本：\n     * 管理成本\n     * 机会成本\n     * 风险准备金\n   - 成本分摊计划（按时间维度）\n   ### 收益预测\n   - 直接收益：\n     * 收入增长\n     * 成本节约\n     * 效率提升\n     * 质量改进\n   - 间接收益：\n     * 品牌价值提升\n     * 客户满意度提升\n     * 市场份额增长\n     * 竞争优势增强\n   - 收益时间分布\n   ### 财务指标预测\n   - 投资回收期（PP）\n   - 净现值（NPV）\n   - 内部收益率（IRR）\n   - 投资回报率（ROI）\n   - 盈亏平衡点分析\n   - 敏感性分析\n   ### 风险评估\n   - 市场风险\n   - 技术风险\n   - 执行风险\n   - 财务风险\n   - 风险应对策略\n3. 使用Markdown格式，结构如下：\n   ## 方案1：方案名称\n   ### 方案概述\n   - 背景和目标\n   - 实施范围\n   - 预期效果\n   ### 实施计划\n   - 具体步骤\n   - 时间安排\n   - 资源需求\n   - 里程碑节点\n   ### ROI预测分析\n   [上述所有ROI预测内容]\n   ### 方案评估\n   - 优势分析\n   - 劣势分析\n   - 机会分析\n   - 威胁分析\n   - 可行性评估\n   [其他方案类似]\n4. 最后给出最优方案推荐，包括：\n   ### 方案对比分析\n   - 各方案关键指标对比表\n   - 优劣势对比分析\n   - 风险对比分析\n   ### 最优方案选择\n   - 选择理由\n   - 实施建议\n   - 预期效果\n   - 风险控制措施\n   - 成功关键因素\n5. 重要数据用**加粗**强调，所有预测过程需要清晰展示\n"""
    try:
        result = call_deepseek(prompt, timeout=300)
        if not result:
            raise Exception("ROI预测生成失败")
        return result
    except Exception as e:
        error_msg = f"ROI预测失败: {str(e)}"
        send_markdown_message(error_msg)
        return error_msg

def handle_execution_intent(text):
    """处理执行意图，支持广告参数生成意图"""
    if any(k in text.lower() for k in ["google ads", "谷歌广告", "广告配置", "生成广告参数", "广告投放参数"]):
        ad_config = handle_ad_config_intent(text)
        if ad_config:
            send_markdown_message("广告参数json已生成，可在侧边栏查看。");
            return {"ad_config_json": ad_config, "report": "广告参数json已生成，可在侧边栏查看。"}
        else:
            error_msg = "广告参数json生成失败，请重试。"
            send_markdown_message(error_msg)
            return {"ad_config_json": None, "report": error_msg}
    prompt = f"""请根据以下执行需求生成具体的执行方案：\n\n需求：{text}\n\n要求：\n1. 如果是自动投放，请说明投放策略和步骤\n2. 如果是其他执行需求，请提供详细的执行方案\n3. 使用Markdown格式\n4. 结构清晰，层次分明\n"""
    try:
        result = call_deepseek(prompt, timeout=300)
        if not result:
            raise Exception("执行方案生成失败")
        send_markdown_message(result)
        return result
    except Exception as e:
        error_msg = f"执行方案生成失败: {str(e)}"
        send_markdown_message(error_msg)
        return error_msg

def handle_content_generation_intent(text, retry_count=0, previous_feedback=None):
    """处理内容生成意图"""
    # 检查是否是小红书内容生成请求
    is_xiaohongshu = any(keyword in text.lower() for keyword in [
        "小红书", "种草", "笔记", "分享", "测评", "推荐", "安利"
    ])
    
    if is_xiaohongshu:
        try:
            # 如果是重试，将上次的问题反馈加入到提示词中
            if retry_count > 0 and previous_feedback:
                text = f"{text}\n\n上次生成的问题：{previous_feedback}"
            
            # 1. 分析场景、角色、性格和行文风格
            style_prompt = f"""请分析以下小红书内容需求，确定场景、角色、性格和行文风格：

内容需求：{text}

请严格按照以下JSON格式返回分析结果，不要包含任何其他内容：
{{
    "scene": {{
        "setting": "具体场景描述",
        "atmosphere": "场景氛围",
        "key_elements": ["关键元素1", "关键元素2"]
    }},
    "role": {{
        "identity": "角色身份",
        "background": "角色背景",
        "expertise": "专业领域"
    }},
    "personality": {{
        "traits": ["性格特征1", "性格特征2"],
        "tone": "语气特点",
        "values": ["价值观1", "价值观2"]
    }},
    "writing_style": {{
        "style": "写作风格",
        "tone": "语气",
        "emotion": "情感基调",
        "keywords": ["关键词1", "关键词2"]
    }}
}}

要求：
1. 场景分析要具体且符合小红书调性
2. 角色设定要真实可信
3. 性格特征要鲜明且讨喜
4. 行文风格要符合目标受众
5. 必须返回完整的JSON对象
6. 不要包含任何注释或说明
7. 确保所有字段都有值
"""
            style_result = call_deepseek(style_prompt, timeout=300)
            if not style_result:
                raise Exception("风格分析失败")
            
            # 清理和验证JSON结果
            try:
                # 移除可能的Markdown代码块标记
                style_result = style_result.strip()
                if style_result.startswith('```json'):
                    style_result = style_result[7:]
                if style_result.endswith('```'):
                    style_result = style_result[:-3]
                style_result = style_result.strip()
                
                # 移除可能的额外内容
                if '```' in style_result:
                    style_result = style_result.split('```')[0].strip()
                
                # 解析JSON
                style_data = json.loads(style_result)
                
                # 验证必要字段
                required_fields = {
                    'scene': ['setting', 'atmosphere', 'key_elements'],
                    'role': ['identity', 'background', 'expertise'],
                    'personality': ['traits', 'tone', 'values'],
                    'writing_style': ['style', 'tone', 'emotion', 'keywords']
                }
                
                for section, fields in required_fields.items():
                    if section not in style_data:
                        raise ValueError(f"缺少{section}部分")
                    for field in fields:
                        if field not in style_data[section]:
                            raise ValueError(f"{section}中缺少{field}字段")
                        if not style_data[section][field]:
                            raise ValueError(f"{section}中的{field}字段为空")
                
                # 提取数据
                scene = style_data['scene']
                role = style_data['role']
                personality = style_data['personality']
                writing_style = style_data['writing_style']
                
            except json.JSONDecodeError as e:
                print(f"JSON解析错误: {e}")
                print(f"原始返回内容: {style_result}")
                raise Exception("风格分析结果格式错误")
            except ValueError as e:
                print(f"数据验证错误: {e}")
                print(f"原始返回内容: {style_result}")
                raise Exception(f"风格分析数据不完整: {str(e)}")
            
            # 2. 生成标题和话题标签
            title_prompt = f"""请根据以下分析结果生成小红书标题和话题标签：

内容需求：{text}

场景分析：
- 场景：{scene['setting']}
- 氛围：{scene['atmosphere']}
- 关键元素：{', '.join(scene['key_elements'])}

角色设定：
- 身份：{role['identity']}
- 背景：{role['background']}
- 专业：{role['expertise']}

性格特征：
- 特点：{', '.join(personality['traits'])}
- 语气：{personality['tone']}
- 价值观：{', '.join(personality['values'])}

写作风格：
- 风格：{writing_style['style']}
- 语气：{writing_style['tone']}
- 情感：{writing_style['emotion']}
- 关键词：{', '.join(writing_style['keywords'])}

要求：
1. 标题要求：
   - 长度15-20字
   - 包含吸引人的关键词
   - 使用emoji表情
   - 突出核心卖点
   - 符合小红书风格

2. 话题标签要求：
   - 只生成一个最相关的话题标签
   - 使用#号标记
   - 标签要简洁有力
   - 要能引起共鸣

请严格按照以下JSON格式返回，不要包含任何其他内容：
{{
    "title": "标题内容",
    "topic": "#话题标签"
}}
"""
            title_result = call_deepseek(title_prompt, timeout=300)
            if not title_result:
                raise Exception("标题生成失败")
            
            try:
                # 清理和验证标题结果
                title_result = title_result.strip()
                if title_result.startswith('```json'):
                    title_result = title_result[7:]
                if title_result.endswith('```'):
                    title_result = title_result[:-3]
                title_result = title_result.strip()
                
                if '```' in title_result:
                    title_result = title_result.split('```')[0].strip()
                
                title_data = json.loads(title_result)
                title = title_data.get('title', '')
                topic = title_data.get('topic', '')
                
                if not title or not topic:
                    raise ValueError("标题或话题标签为空")
                
                if not topic.startswith('#'):
                    topic = f"#{topic}"
                
            except json.JSONDecodeError as e:
                print(f"标题JSON解析错误: {e}")
                print(f"原始返回内容: {title_result}")
                raise Exception("标题结果格式错误")
            except ValueError as e:
                print(f"标题数据验证错误: {e}")
                print(f"原始返回内容: {title_result}")
                raise Exception(f"标题数据不完整: {str(e)}")
            
            # 3. 生成正文
            content_prompt = f"""请根据以下分析结果生成小红书正文内容：

标题：{title}
话题：{topic}
需求：{text}

场景分析：
- 场景：{scene['setting']}
- 氛围：{scene['atmosphere']}
- 关键元素：{', '.join(scene['key_elements'])}

角色设定：
- 身份：{role['identity']}
- 背景：{role['background']}
- 专业：{role['expertise']}

性格特征：
- 特点：{', '.join(personality['traits'])}
- 语气：{personality['tone']}
- 价值观：{', '.join(personality['values'])}

写作风格：
- 风格：{writing_style['style']}
- 语气：{writing_style['tone']}
- 情感：{writing_style['emotion']}
- 关键词：{', '.join(writing_style['keywords'])}

要求：
1. 内容结构：
   - 开头吸引人，用emoji表情
   - 分段清晰，每段都有重点
   - 包含个人真实感受
   - 突出核心卖点
   - 结尾有互动引导

2. 写作风格：
   - 严格遵循设定的角色性格
   - 保持一致的语气和情感
   - 使用emoji表情
   - 重点内容加粗
   - 适当使用感叹号

3. 内容要求：
   - 字数500-800字
   - 直接输出完整笔记内容
   - 不要包含任何建议或分析
   - 确保内容真实可信
   - 突出个人特色
"""
            content = call_deepseek(content_prompt, timeout=300)
            if not content:
                raise Exception("正文生成失败")
            
            # 整合最终内容
            final_content = f"""# {title}

{content}

{topic}"""
            
            # 发送到企业微信
            send_markdown_message(final_content)
            return final_content
            
        except Exception as e:
            error_msg = f"小红书内容生成失败: {str(e)}"
            send_markdown_message(error_msg)
            return error_msg
    
    else:
        # 其他内容生成使用原有逻辑
        has_sufficient_guidance = any(keyword in text.lower() for keyword in [
            "格式", "结构", "要求", "风格", "markdown", "标题", "段落", "列表",
            "加粗", "引用", "代码块", "概述", "要点", "建议"
        ])
        
        if has_sufficient_guidance:
            prompt = text
        else:
            prompt = f"""你是一个专业的内容创作专家，请根据以下需求生成高质量的内容：

需求：{text}

要求：
1. 根据内容类型选择合适的风格和格式
2. 确保内容专业、准确、有吸引力
3. 使用Markdown格式
4. 结构清晰，层次分明
5. 包含必要的标题、段落和列表
6. 重要内容用**加粗**强调
7. 适当使用引用和代码块
8. 确保内容符合目标受众的需求

请按照以下结构组织内容：
1. 内容概述
2. 主要内容
3. 关键要点
4. 行动建议（如适用）
"""
        
        try:
            result = call_deepseek(prompt, timeout=300)
            
            if result is None:
                error_msg = "内容生成失败，请重试"
                send_markdown_message(error_msg)
                return error_msg
                
            send_markdown_message(result)
            return result
        except Exception as e:
            error_msg = f"内容生成失败: {str(e)}"
            send_markdown_message(error_msg)
            return error_msg

def handle_geo_marketing_intent(text):
    """处理GEO营销相关意图，包括收录统计和收录解决方案生成"""
    print("handle_geo_marketing_intent")
    print("--text--")
    print(text)
    # 判断是否为收录统计
    if "收录统计" in text or "收录情况" in text or "Google收录" in text:
        # 提取用户输入的真实query（双引号内内容）
        import re
        match = re.search(r'"([^"]+)"', text)
        print("--match--")
        print(match)
        if match:
            query = match.group(1)
        else:
            # 兼容无双引号时，去除常见前缀
            #query = text.replace("查询", "").replace("搜索", "").replace("收录统计", "").replace("收录情况", "").replace("Google收录", "").strip()
            query = text
            if not query:
                query = "一键式Google广告投放"  # 默认查询词
        print("--query--")
        print(query)
        # 使用kimi联网搜索
        response = kimi_web_search(query)
        print("--response--")
        print(response)
        if response:
            ads_list, organic_list = extract_search_data(response, TARGET_BRAND)
            print("--ads_list--")
            print(ads_list)
            print("--organic_list--")
            print(organic_list)
            total_found = sum(ad['brand_mentioned'] for ad in ads_list + organic_list)
            
            # 生成收录统计报告
            report = f"## Google收录统计报告\n\n"
            report += f"**查询关键词**: {query}\n"
            report += f"**目标品牌**: {TARGET_BRAND}\n"
            report += f"**收录情况**: 在{len(ads_list) + len(organic_list)}条搜索结果中，目标品牌出现{total_found}次\n\n"
            
            if ads_list:
                report += "### 广告位收录情况\n"
                for ad in ads_list:
                    status = "✅ 提及" if ad['brand_mentioned'] else "❌ 未提及"
                    report += f"- **排名 #{ad['position']}** | {status} | {ad['content_snippet']}\n"
                    if ad['source_links']:
                        report += f"  - 链接: {ad['source_links'][0]}\n"
                report += "\n"
            
            if organic_list:
                report += "### 自然搜索结果收录情况\n"
                for org in organic_list:
                    status = "✅ 提及" if org['brand_mentioned'] else "❌ 未提及"
                    report += f"- **排名 #{org['position']}** | {status} | {org['content_snippet']}\n"
                    if org['source_links']:
                        report += f"  - 链接: {org['source_links'][0]}\n"
                report += "\n"
            
            # 相关性分析
            report += "### 相关性分析\n"
            if total_found == 0:
                report += "❌ **相关性较低**: 目标品牌在当前搜索结果中未被提及，建议优化品牌曝光策略。\n"
            elif total_found < 3:
                report += "⚠️ **相关性一般**: 目标品牌提及次数较少，建议加强品牌推广。\n"
            else:
                report += "✅ **相关性较高**: 目标品牌在搜索结果中表现良好。\n"
            
            return report
        else:
            return "收录统计查询失败，请稍后重试。"
    
    # 判断是否为收录解决方案生成
    if "收录解决方案" in text or "收录优化" in text or "提升收录" in text or "SEO方案" in text:
        # 提取用户输入的查询关键词
        query = text
        query = query.replace("收录解决方案", "").replace("收录优化", "").replace("提升收录", "").replace("SEO方案", "").strip()
        if not query:
            query = "一键式Google广告投放"
        
        # 先进行收录统计
        response = kimi_web_search(query)
        if response:
            ads_list, organic_list = extract_search_data(response, TARGET_BRAND)
            total_found = sum(ad['brand_mentioned'] for ad in ads_list + organic_list)
            
            # 生成解决方案
            solution = f"## 收录优化解决方案\n\n"
            solution += f"**查询关键词**: {query}\n"
            solution += f"**目标品牌**: {TARGET_BRAND}\n"
            solution += f"**当前收录**: {total_found}次提及\n\n"
            
            solution += "### 优化建议\n"
            solution += generate_solution(TARGET_BRAND, total_found)
            
            # 针对性的优化策略
            if total_found == 0:
                solution += "\n### 紧急优化策略\n"
                solution += "1. **立即投放Google Ads品牌词广告**\n"
                solution += "2. **优化官网SEO，确保品牌词收录**\n"
                solution += "3. **发布品牌相关高质量内容**\n"
                solution += "4. **建立品牌外链和社交媒体曝光**\n"
            elif total_found < 3:
                solution += "\n### 提升策略\n"
                solution += "1. **扩大关键词覆盖范围**\n"
                solution += "2. **优化现有内容质量**\n"
                solution += "3. **增加长尾关键词布局**\n"
                solution += "4. **监控竞争对手动态**\n"
            
            return solution
        else:
            return "收录解决方案生成失败，请稍后重试。"
    
    return "暂不支持的GEO营销意图。"

@app.route('/chatss', methods=['POST'])
@login_required
def chat_handlerss():
    try:
        # 检查用户是否存在
        if not current_user or not current_user.is_authenticated:
            return jsonify({'error': '请先登录'}), 401

        # 获取用户实例（同时获取SQLite和MySQL中的用户）
        user = User.query.get(current_user.id)
        mysql_user = MySQLUser.query.get(current_user.id)
        
        if not user or not mysql_user:
            return jsonify({'error': '用户不存在'}), 404

        # 获取或创建用户每日使用记录
        today = datetime.now().date()
        daily_usage = UserDailyUsage.get_or_create(current_user.id, today)
        
        # 检查用户是否达到每日查询限制
        if not daily_usage.can_make_query():
            return jsonify({'error': '已达到今日查询限制（100次）'}), 429

        data = request.json
        message = data.get('message', '')
        file_id = data.get('fileId')
        
        if not message:
            return jsonify({'error': '消息不能为空'}), 400

        # 获取文件数据（如果提供了文件ID）
        file_data = None
        if file_id and file_id in file_storage:
            file_data = file_storage[file_id]

        # 识别用户意图
        intent_result = recognize_intent(message)
        if not intent_result:
            return jsonify({'error': '无法识别用户意图'}), 400

        def generate():
            try:
                # 根据意图处理请求
                result = None
                # IP位置查询意图处理
                if "IP位置查询意图" in intent_result.get("intents", []):
                    # 尝试从用户输入中提取IP地址
                    import re
                    ip_match = re.search(r"(\d{1,3}(?:\.\d{1,3}){3})", message)
                    ip = ip_match.group(1) if ip_match else "8.8.8.8"
                    result = call_mcp_tool('get_ip_info', {'ip': ip})
                elif any(x in intent_result.get("intents", []) for x in ["GEO营销-收录统计", "GEO营销-收录解决方案生成"]):
                    result = handle_geo_marketing_intent(message)
                elif "分析意图" in intent_result.get("intents", []):
                    result = handle_analysis_intent(message, file_data)
                elif "ROI预测意图" in intent_result.get("intents", []):
                    result = handle_roi_intent(message, file_data)
                elif "执行意图" in intent_result.get("intents", []):
                    result = handle_execution_intent(message)
                elif "内容生成意图" in intent_result.get("intents", []):
                    result = handle_content_generation_intent(message)
                else:
                    result = "抱歉，我无法理解您的意图。请尝试重新描述您的问题。"
                if not result:
                    result = "抱歉，处理您的请求时出现错误。请稍后重试。"

                # 在应用上下文中执行数据库操作
                with app.app_context():
                    # 增加用户查询计数（同时更新所有相关表）
                    if not user.increment_queries() or not mysql_user.increment_queries() or not daily_usage.increment_usage():
                        yield f"data: {json.dumps({'content': '更新查询次数失败，请重试'})}\n\n"
                        yield "data: [DONE]\n\n"
                        return

                    # 保存聊天记录
                    try:
                        import json
                        model_response_str = result
                        if not isinstance(model_response_str, str):
                            try:
                                model_response_str = json.dumps(model_response_str, ensure_ascii=False)
                            except Exception:
                                model_response_str = str(model_response_str)
                        chat_record = ChatRecord(
                            user_id=user.id,
                            question=message,
                            intent_result=json.dumps(intent_result),
                            model_response=model_response_str,
                            file_id=file_id
                        )
                        db.session.add(chat_record)
                        
                        # 保存问答记录
                        qa_record = QARecord(
                            conversation_id=QARecord.generate_conversation_id(),
                            user_id=user.id,
                            question=message,
                            answer=result,
                            user_tag=intent_result.get("intents", [""])[0],  # 使用第一个意图作为标签
                            answer_quality="待评估",  # 初始状态
                            meta_data={
                                "intent_result": intent_result,
                                "file_id": file_id,
                                "has_file": bool(file_data)
                            }
                        )
                        db.session.add(qa_record)
                        db.session.commit()
                    except Exception as e:
                        print(f"Error saving records: {str(e)}")
                        db.session.rollback()

                # 生成报告标题
                report_title = generate_report_title(message, result)
                
                # 生成PPT报告
                try:
                    pptx_path = text_to_pptxs(result, report_title)
                    if pptx_path:
                        upload_file_to_wecom(pptx_path, webhook_url)
                except Exception as e:
                    print(f"生成PPT报告失败: {str(e)}")
                
                # 流式返回结果
                yield f"data: {json.dumps({'content': result})}\n\n"
                yield "data: [DONE]\n\n"
                
                # 广告参数生成意图处理
                if "广告参数生成意图" in intent_result.get("intents", []):
                    ad_config = handle_ad_config_intent(message)
                    result = "广告参数json已生成，可在侧边栏查看。"
                    print("--ad_config--")
                    print(ad_config)
                    print(type(ad_config))
                    print("--result--")
                    print(result)
                    yield f"data: {{\"ad_config_json\": {json.dumps(ad_config, ensure_ascii=False)}, \"content\": result}}\n\n"
                    yield "data: [DONE]\n\n"
                    return

            except Exception as e:
                print(f"Error in generate: {str(e)}")
                yield f"data: {json.dumps({'content': '抱歉，处理您的请求时出现错误。请稍后重试。'})}\n\n"
                yield "data: [DONE]\n\n"

        return Response(generate(), mimetype='text/event-stream')
            
    except Exception as e:
        print(f"Error in chat_handlerss: {str(e)}")
        return jsonify({'error': '处理请求时出现错误'}), 500

import tempfile
import os
from pptx import Presentation
from pptx.util import Pt

def generate_report_title(question, answer):
    """生成6字标题"""
    prompt = f"""请根据以下问题和答案生成一个准确、吸引人的标题,标题不超过6个字：

问题：{question}
答案：{answer}

要求：
1. 要概括内容核心
2. 要简洁有力
3. 要正式专业
4. 要表明具体方向

例子：
1、询盘AI策略优化_20250530
"""

    try:
        response = client.chat.completions.create(
            model="glm-4",
            messages=[{"role": "user", "content": prompt}],
            stream=False
        )
        title = response.choices[0].message.content.strip()

        print("--title result--")
        print(question)
        print(answer)
        print(title)
        # 确保返回的是6个汉字
        #if len(title) != 6:
        #    return "数据分析报告"
        return title
    except:
        return "数据分析报告"

def structure_content_for_ppt(text, title):
    """使用GLM-4对内容进行结构化，生成适合PPT展示的格式"""
    prompt = f"""请将以下内容重新组织为适合PPT展示的格式：

原始内容：
{text}

要求：
1. 保持专业性和逻辑性
2. 每个段落应该是一个独立的要点
3. 使用简洁的语言
4. 确保内容层次分明
5. 每个段落都应该有一个清晰的主题,标题不超过6个字
6. 使用"#"标记每个段落的标题
7. 标题要简短有力
8. 内容要点化，便于PPT展示

请按照以下格式输出：
# 标题1
内容1

# 标题2
内容2

...以此类推
"""
    
    try:
        response = client.chat.completions.create(
            model="glm-4",
            messages=[{"role": "user", "content": prompt}],
            stream=False
        )
        structured_content = response.choices[0].message.content

        print("--structured_content--")
        print(structured_content)
        return structured_content
    except Exception as e:
        print(f"Error structuring content: {e}")
        return text

def text_to_pptx(text, title="AI分析报告", custom_prefix=None):
    """使用AiPPT API创建PPT"""
    try:
        # 初始化AiPPT客户端
        ai_ppt = AIPPT(AIPPT_APP_ID, AIPPT_API_SECRET, text, AIPPT_TEMPLATE_ID)
        
        # 创建PPT生成任务
        task_id = ai_ppt.create_task()
        if not task_id:
            print("创建PPT任务失败，使用备用方法")
            return text_to_pptxs(text, title)
        
        # 获取PPT下载链接
        try:
            ppt_url = ai_ppt.get_result(task_id)
            if not ppt_url:
                print("获取PPT下载链接失败，使用备用方法")
                return text_to_pptxs(text, title)
        except json.JSONDecodeError as e:
            print(f"JSON解析错误: {e}，使用备用方法")
            return text_to_pptxs(text, title)
        except Exception as e:
            print(f"获取PPT下载链接时发生错误: {e}，使用备用方法")
            return text_to_pptxs(text, title)
        
        # 下载PPT文件
        try:
            response = requests.get(ppt_url)
            if response.status_code != 200:
                print(f"下载PPT失败，状态码: {response.status_code}，使用备用方法")
                return text_to_pptxs(text, title)
        except Exception as e:
            print(f"下载PPT时发生错误: {e}，使用备用方法")
            return text_to_pptxs(text, title)
        
        # 保存到临时文件
        try:
            if custom_prefix:
                fd, tmp_path = tempfile.mkstemp(suffix='.pptx', prefix=custom_prefix + '_')
            else:
                fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
            
            with open(tmp_path, 'wb') as f:
                f.write(response.content)
            
            os.close(fd)
            return tmp_path
        except Exception as e:
            print(f"保存PPT文件失败: {e}，使用备用方法")
            return text_to_pptxs(text, title)
        
    except Exception as e:
        print(f"Error generating PPT: {e}，使用备用方法")
        # 如果API调用失败，回退到原来的PPT生成方法
        return text_to_pptxs(text, title)

def text_to_pptxs(text, title="AI分析报告"):
    prs = Presentation()
    # 标题页
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "由AI自动生成"

    # 内容页
    slide_layout = prs.slide_layouts[1]
    for idx, para in enumerate(text.split('\n\n')):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"第{idx+1}部分"
        content = slide.placeholders[1]
        content.text = para.strip()
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)

    # 保存到临时文件
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def upload_file_to_wecom(filepath, webhook_url):
    # 获取上传url
    upload_url = webhook_url.replace('/send?', '/upload_media?') + '&type=file'
    
    with open(filepath, 'rb') as f:
        print("--filepath--")
        print(filepath)
        print(f)
        files = {'media': f}
        resp = requests.post(upload_url, files=files)
    media_id = resp.json().get('media_id')
    return media_id

def send_file_to_wecom(media_id, webhook_url):
    data = {
        "msgtype": "file",
        "file": {
            "media_id": media_id
        }
    }
    resp = requests.post(webhook_url, json=data)
    return resp.json()

@app.route('/wecom_webhook', methods=['POST'])
def wecom_webhook():
    """处理企业微信机器人的webhook请求"""
    try:
        # 获取请求数据
        data = request.json
        print("--wecom webhook data--")
        print(data)
        
        # 提取消息内容
        msg_type = data.get('msgtype')
        if msg_type == 'text':
            content = data.get('text', {}).get('content', '').strip()
            # 移除@机器人的标记
            content = content.replace('@MarketManus', '').strip()
            
            # 检查是否有文件上传
            file_id = None
            if 'file' in data:
                file_id = data['file'].get('file_id')
            
            # 调用现有的聊天处理逻辑
            chat_data = {
                'message': content,
                'fileId': file_id
            }
            
            # 使用现有的chat_handlerss函数处理请求
            response = chat_handlerss()
            return response
            
        else:
            return jsonify({'error': '不支持的消息类型'}), 400
            
    except Exception as e:
        print(f"Error processing WeChat Work webhook: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/qa/quality/<conversation_id>', methods=['PUT'])
@login_required
def update_qa_quality(conversation_id):
    """更新问答质量评估"""
    try:
        data = request.json
        quality = data.get('quality')
        
        if not quality:
            return jsonify({'error': '质量评估不能为空'}), 400
            
        # 查找对应的问答记录
        qa_record = QARecord.query.filter_by(conversation_id=conversation_id).first()
        if not qa_record:
            return jsonify({'error': '问答记录不存在'}), 404
            
        # 更新质量评估
        qa_record.answer_quality = quality
        db.session.commit()
        
        return jsonify({
            'message': '质量评估已更新',
            'qa_record': qa_record.to_dict()
        }), 200
        
    except Exception as e:
        print(f"Error updating QA quality: {e}")
        return jsonify({'error': '更新质量评估失败'}), 500

# 创建数据库表
with app.app_context():
    # 创建所有数据库表
    db.create_all()

# Kimi API参数
KIMI_API_KEY = "sk-8mUEjnAhKkQaD5c1BVodRr7mbdJeyRZHNk7HYozxUOmdBOaa"  # 替换为实际 API Key
SEARCH_KEYWORD = "一键式Google广告投放"
TARGET_BRAND = "云链灵智MarketManus智能投放引擎"

def kimi_web_search(query: str, retry=2):
    """调用 Kimi 官方 API 执行联网搜索，query为用户输入的真实搜索词，失败自动重试"""
    import os
    url = "https://api.moonshot.cn/v1/chat/completions"
    #api_key = os.environ.get("MOONSHOT_API_KEY")
    api_key = KIMI_API_KEY

    if not api_key:
        raise ValueError("请设置环境变量 MOONSHOT_API_KEY")
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "moonshot-v1-128k",
        "messages": [
            {"role": "system", "content": "你是 Kimi。"},
            {"role": "user", "content": query}
        ],
        "tools": [
            {
                "type": "builtin_function",
                "function": {
                    "name": "$web_search"
                }
            }
        ],
        "temperature": 0.3
    }
    for i in range(retry):
        try:
            response = requests.post(url, json=payload, headers=headers, timeout=60)
            if response.status_code != 200:
                print(f"API 请求失败: {response.status_code} {response.text}，重试 {i+1}/{retry}")
                continue
            return response.json()
        except Exception as e:
            print(f"API 请求异常: {e}，重试 {i+1}/{retry}")
    return None

def extract_search_data(response: dict, target: str):
    """从 Kimi 返回内容中提取广告和自然结果信息"""
    content = response['choices'][0]['message']['content']
    # 解析广告和自然结果区块
    ad_pattern = re.compile(r'\d+\.\s*广告[\s\S]*?(?=\d+\.\s*(广告|自然结果)|$)')
    organic_pattern = re.compile(r'\d+\.\s*自然结果[\s\S]*?(?=\d+\.\s*(广告|自然结果)|$)')
    ad_sections = ad_pattern.findall(content)
    organic_sections = organic_pattern.findall(content)
    # 提取引用链接
    ref_links = re.findall(r'\[(\d+)\]:\s*(http\S+)', content)
    ref_dict = {num: link for num, link in ref_links}
    def parse_section(sections, result_type):
        result = []
        for idx, sec in enumerate(sections):
            brand_mentioned = target in sec
            ad_refs = re.findall(r'\[(\d+)\]', sec)
            links = [ref_dict.get(ref) for ref in ad_refs if ref in ref_dict]
            result.append({
                "type": result_type,
                "position": idx + 1,
                "content_snippet": sec.split('\n')[0][:100],
                "brand_mentioned": brand_mentioned,
                "source_links": links
            })
        return result
    ads_data = parse_section(ad_sections, "广告")
    organic_data = parse_section(organic_sections, "自然结果")
    return ads_data, organic_data

def generate_solution(brand, found_count):
    """根据收录情况生成SEO/SEM优化建议"""
    if found_count == 0:
        return f'未在Google广告/自然结果中发现"{brand}"的收录。建议：\n1. 优化品牌官网SEO，提升品牌词和长尾词收录。\n2. 增加高质量外链，提升权重。\n3. 结合Google Ads投放品牌词广告，提升曝光。\n4. 检查robots.txt和站点地图，确保无技术屏蔽。\n5. 丰富内容，覆盖更多用户搜索意图。\n'
    elif found_count < 3:
        return f'"{brand}"收录较少。建议：\n1. 持续优化SEO，关注排名提升。\n2. 增加内容多样性，覆盖更多长尾词。\n3. 监控广告投放效果，及时调整关键词和出价。\n'
    else:
        return f'"{brand}"收录表现良好，请持续优化并关注竞争对手动态。'
        

# 可在GEO营销意图处理逻辑中调用如下：
# response = kimi_web_search(SEARCH_KEYWORD)
# if response:
#     ads_list, organic_list = extract_search_data(response, TARGET_BRAND)
#     total_found = sum(ad['brand_mentioned'] for ad in ads_list + organic_list)
#     print(f"目标品牌 '{TARGET_BRAND}' 出现在 {total_found} 条结果中：")
#     for ad in ads_list:
#         print(f"[广告] 排名 #{ad['position']} | 提及: {'是' if ad['brand_mentioned'] else '否'} | 链接: {ad['source_links'][0] if ad['source_links'] else '无'}")
#     for org in organic_list:
#         print(f"[自然] 排名 #{org['position']} | 提及: {'是' if org['brand_mentioned'] else '否'} | 链接: {org['source_links'][0] if org['source_links'] else '无'}")
#     print("---\n【收录优化建议】")
#     print(generate_solution(TARGET_BRAND, total_found))
# else:
#     print("搜索失败，请检查 API 或网络设置")

def handle_ad_config_intent(text, user_params=None):
    """生成广告参数json，one-shot content.json，融合用户参数，并保存到本地和远程服务器"""
    import json
    import uuid
    import os
    import subprocess
    # 读取one-shot示例
    print("--start handle_ad_config_intent--")
    with open('content.json', 'r', encoding='utf-8-sig') as f:
        one_shot = json.load(f)
    print("--one_shot--")
    print(one_shot)
    # 构造大模型提示
    prompt = f"""你是一个专业的广告投放参数生成助手。请根据用户需求生成结构化的广告参数json，格式严格参考以下示例：\n\n【示例】\n{json.dumps(one_shot, ensure_ascii=False, indent=2)}\n\n【用户需求】\n{text}\n\n【要求】\n1. 输出必须是完整的json对象，字段结构与示例一致。\n2. 如果用户指定了部分参数（如预算、地域、关键词等），请优先使用用户指定，未指定的部分参考示例自动补全。\n3. 只返回json，不要有多余解释。"""
    result = call_deepseek(prompt, timeout=300)
    # 尝试解析json，自动去除BOM
    try:
        result = result.strip()
        if result.startswith('```json'):
            result = result[7:]
        if result.endswith('```'):
            result = result[:-3]
        result = result.strip()
        # 处理BOM头
        if result and ord(result[0]) == 0xfeff:
            result = result[1:]
        ad_config = json.loads(result)
        # 保存到本地唯一json文件
        file_id = str(uuid.uuid4())
        file_name = f"ad_config_{file_id}.json"
        file_path = os.path.join("./", file_name)
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(ad_config, f, ensure_ascii=False, indent=2)
        # 通过scp传输到远程服务器（占位符信息，后期替换）
        remote_ip = "10.40.228.183"  # TODO: 替换为实际ip
        remote_user = "root"  # TODO: 替换为实际用户名
        remote_path = "/root/ad_config/"  # TODO: 替换为实际路径
        try:
            scp_cmd = [
                "scp",
                file_path,
                f"{remote_user}@{remote_ip}:{remote_path}{file_name}"
            ]
            subprocess.run(scp_cmd, check=True)
            print(f"ad_config_json已通过scp传输到远程: {remote_ip}")
        except Exception as e:
            print(f"scp传输失败: {e}")
        return ad_config
    except Exception as e:
        print(f"广告参数json解析失败: {e}")
        return None

# MCP工具调用函数
def call_mcp_tool(tool_name, params, mcp_url="http://10.107.176.206:8080"):
    url = f"{mcp_url}/tool/{tool_name}"
    print(f"MCP调用URL: {url}")
    try:
        resp = requests.post(url, json=params, timeout=10)
        resp.raise_for_status()
        data = resp.json()
        if 'result' in data:
            return data['result']
        else:
            raise Exception(data.get('error', 'Unknown MCP tool error'))
    except Exception as e:
        return f"MCP调用失败: {str(e)}"

@app.route('/ipinfo')
def get_ipinfo():
    ip = request.args.get('ip', '8.8.8.8')
    res = call_mcp_tool('get_ip_info', {'ip': ip})
    return jsonify({"data": res})

if __name__ == '__main__':
    app.run(host="0.0.0.0", port=8080) 